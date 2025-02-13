import fitz  # PyMuPDF is imported as fitz
import re
import docx
from pptx import Presentation
import openpyxl
import msoffcrypto

def mask_pii_in_text(text):
    if text is None:
        return ""

    # Mask email addresses (show first 2 characters)
    text = re.sub(r'\b([A-Za-z0-9._%+-]{2})[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,7}\b', r'\1***@***.com', text)

    # Mask phone numbers (global format)
    text = re.sub(r'\b(\+?\d{1,4}[-.\s]??\(?\d{1,4}\)?[-.\s]??\d{1,4}[-.\s]??\d{1,9})\b', 'XXX-XXX-XXXX', text)

    # Mask Aadhaar numbers (12-digit format, no leading 0 or 1, with spaces)
    text = re.sub(r'^[2-9]{1}[0-9]{3}\s[0-9]{4}\s[0-9]{4}$', 'XXXX XXXX XXXX', text)

    # Mask PAN card numbers (10 characters, 5 letters, 4 digits, 1 letter)
    text = re.sub(r'[A-Z]{5}[0-9]{4}[A-Z]{1}', 'XXXXXXXXXX', text)

    # Mask Driving License numbers (specific format as per Indian DL, keeping first 2 characters visible)
    text = re.sub(r'^([A-Z]{2})[0-9]{2}( |-|)[0-9]{4}(19|20)[0-9]{2}[0-9]{7}$', r'\1XX XXXX XXXX XXXX', text)

    return text

def modify_and_encrypt_pdf(input_stream, output_stream, password):
    try:
        # Create a new PDF document from the stream
        doc = fitz.open(stream=input_stream.read(), filetype="pdf")
        
        # Process each page
        for page in doc:
            # Process text on each page
            text = page.get_text()
            # Mask PII in text if needed
            masked_text = mask_pii_in_text(text)
            if text != masked_text:
                # If text was modified, update the page
                page.clean_contents()
                page.insert_text((50, 50), masked_text)  # You may need to adjust coordinates
        
        # Save with encryption
        doc.save(
            output_stream,
            encryption=fitz.PDF_ENCRYPT_AES_256,  # Use AES 256-bit encryption
            user_pw=password,  # User password
            owner_pw=password,  # Owner password
            permissions=int(
                fitz.PDF_PERM_PRINT |  # Allow printing
                fitz.PDF_PERM_COPY    # Allow copying text
            )
        )
        doc.close()
        return True
    except Exception as e:
        print(f"Error in modify_and_encrypt_pdf: {str(e)}")
        raise

def modify_docx(input_stream, output_stream):
    doc = docx.Document(input_stream)

    # Iterate through each paragraph and replace PII
    for paragraph in doc.paragraphs:
        original_text = paragraph.text
        masked_text = mask_pii_in_text(original_text)
        if original_text != masked_text:
            paragraph.text = masked_text

    # Iterate through each table and replace PII
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                original_text = cell.text
                masked_text = mask_pii_in_text(original_text)
                if original_text != masked_text:
                    cell.text = masked_text

    # Save the modified document
    doc.save(output_stream)

def mask_pptx_file(input_stream, output_stream):
    prs = Presentation(input_stream)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        masked_text = mask_pii_in_text(original_text)
                        if original_text != masked_text:
                            run.text = masked_text

    prs.save(output_stream)

def mask_excel_file(input_stream, output_stream):
    wb = openpyxl.load_workbook(input_stream)

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    original_text = cell.value
                    masked_text = mask_pii_in_text(original_text)
                    if original_text != masked_text:
                        cell.value = masked_text

    wb.save(output_stream)

def encrypt_excel_file(input_stream):
    office_file = msoffcrypto.OfficeFile(input_stream)
    office_file.load_key(password="securepassword")
    office_file.encrypt("securepassword", input_stream)