import fitz  # PyMuPDF
import re
import docx
from pptx import Presentation
import openpyxl
import msoffcrypto

def mask_pii_in_text(text):
    if text is None:
        return ""

    # Mask email addresses (show first 2 characters)
    text = re.sub(r'\b([A-Za-z0-9._%+-]{2})[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,7}\b', r'\1***@***.com', text)

    # Mask phone numbers (global format)
    text = re.sub(r'\b(\+?\d{1,4}[-.\s]??\(?\d{1,4}\)?[-.\s]??\d{1,4}[-.\s]??\d{1,4}[-.\s]??\d{1,9})\b', 'XXX-XXX-XXXX', text)

    # Mask Aadhaar numbers (12-digit format, no leading 0 or 1, with spaces)
    text = re.sub(r'^[2-9]{1}[0-9]{3}\s[0-9]{4}\s[0-9]{4}$', 'XXXX XXXX XXXX', text)

    # Mask PAN card numbers (10 characters, 5 letters, 4 digits, 1 letter)
    text = re.sub(r'[A-Z]{5}[0-9]{4}[A-Z]{1}', 'XXXXXXXXXX', text)

    # Mask Driving License numbers (specific format as per Indian DL, keeping first 2 characters visible)
    text = re.sub(r'^([A-Z]{2})[0-9]{2}( |-|)[0-9]{4}(19|20)[0-9]{2}[0-9]{7}$', r'\1XX XXXX XXXX XXXX', text)

    return text

def modify_and_encrypt_pdf(input_pdf_path, output_pdf_path, encrypted_pdf_path, password):
    # Open the document
    doc = fitz.open(input_pdf_path)

    # Mask PII and redact the content
    for page in doc:
        blocks = page.get_text("blocks")
        for block in blocks:
            block_text = block[4]
            masked_text = mask_pii_in_text(block_text)

            if masked_text != block_text:  # If there's any change, replace it
                # Redact the original block (visually hide it)
                page.add_redact_annot(block[:4], fill=(1, 1, 1))  # White color fill for redaction
                page.apply_redactions()

                # Insert the masked text at the same location
                page.insert_text((block[0], block[1]), masked_text, fontsize=12)  # Adjust fontsize as needed

    # Save the modified PDF
    doc.save(output_pdf_path)

    # Encrypt the modified PDF with a single password
    doc = fitz.open(output_pdf_path)
    doc.save(encrypted_pdf_path, encryption=fitz.PDF_ENCRYPT_AES_256, owner_pw=password, user_pw=password)
    doc.close()

def modify_docx(input_docx_path, output_docx_path):
    # Load the document
    doc = docx.Document(input_docx_path)

    # Iterate through each paragraph and replace PII
    for paragraph in doc.paragraphs:
        original_text = paragraph.text
        masked_text = mask_pii_in_text(original_text)
        if original_text != masked_text:
            paragraph.text = masked_text

    # Iterate through each table and replace PII
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                original_text = cell.text
                masked_text = mask_pii_in_text(original_text)
                if original_text != masked_text:
                    cell.text = masked_text

    # Save the modified document
    doc.save(output_docx_path)

def mask_pptx_file(input_pptx_path, output_pptx_path):
    prs = Presentation(input_pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text
                        masked_text = mask_pii_in_text(original_text)
                        if original_text != masked_text:
                            run.text = masked_text

    prs.save(output_pptx_path)

def mask_excel_file(input_excel_path, output_excel_path):
    wb = openpyxl.load_workbook(input_excel_path)

    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    original_text = cell.value
                    masked_text = mask_pii_in_text(original_text)
                    if original_text != masked_text:
                        cell.value = masked_text

    wb.save(output_excel_path)

def encrypt_excel_file(input_excel_path, output_excel_path, password):
    with open(input_excel_path, 'rb') as file:
        office_file = msoffcrypto.OfficeFile(file)
        office_file.load_key(password=password)

        with open(output_excel_path, 'wb') as encrypted_file:
            office_file.encrypt(password, encrypted_file)